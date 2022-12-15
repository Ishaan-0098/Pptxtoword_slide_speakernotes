#CREATOR = ISHAAN##
from tkinter import *
import argparse
import csv
import json
from docx.shared import Mm
from docx.shared import Pt
import random
from docx import Document
from docx.shared import Inches
import os
import time
from PIL import Image
from PIL import ImageDraw
from PIL import ImageFont
from functools import lru_cache
import collections 
import collections.abc
from pptx import Presentation
import re
from tkinter import *
from tkinter import filedialog
import sys
from docxcompose.composer import Composer
from docx import Document as Document_compose
from pptx_tools.utils import save_pptx_as_png
import os
import comtypes.client
import time


top = Tk()  
top.geometry("600x300") 
top.title("Doc Creator By Ishaan")
#top.iconbitmap(r'D:\icons\favicon.ico')
#top.configure(bg='pink')

def resource_path2(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller """
    try:
        # PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")
 
    return os.path.join(base_path, relative_path)
def resource_path0(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller """
    base_path = getattr(
        sys,
        '_MEIPASS',
        os.path.dirname(os.path.abspath(__file__)))
    return os.path.join(base_path, relative_path)
 
 
def resource_path(relative):
    return os.path.join(
        os.environ.get(
            "_MEIPASS2",
            os.path.abspath(".")
        ),
        relative
    )
 
    return os.path.join(base_path, relative_path)

def give_path():
    temppath=(r'C:\Users\Tnluser\Desktop\Doc Creator By Ishaan\template (1).docx')

    temppath2=str(my_str3.get())
    pathy = resource_path2(r"template (1).docx")

    
    #pptfile = str(path_value.get())
    pptfile_p= str(my_str1.get())
    pptfile_res=str(my_str2.get())
    pptfile_p=pptfile_p.replace('/','\\')
    pptfile_res=pptfile_res.replace('/','\\')
    #print(pptfile_p)
    #print(pptfile_res)
    #for images
    
    dire = pptfile_res+"\\pngs"
    #print(dire)
    listy=[]

        # use full path for pptx_filename
    listy.append(save_pptx_as_png(dire ,pptfile_p, overwrite_folder=True))
    #print(listy)
    os.chdir(dire)
    for f in os.listdir():
        file_name,file_ext = os.path.splitext(f)
        f_title, f_num = file_name.split('e')
        f_title = f_title.strip()
        f_num = f_num.strip().zfill(2)
        #print(f_num)
        new_name='{}{}'.format(f_num,file_ext)
        
        #print(new_name)
    #sorted(os.listdir())

        os.rename(f,new_name)
        
    #for i in range(len(list)):
     #   print(list[i])


    #for picture paths to automate later
    


    SUPPORTED_IMAGES = ['.jpg', '.png', '.jpeg']

    MAX_IMAGES_PER_DIR = 1000

   
    all_images = []
    for root, _, files in os.walk(dire):
        files = [os.path.join(root, f) for f in filter(lambda x: os.path.splitext(x)[-1].lower() in SUPPORTED_IMAGES, files)]
        all_images.extend(files)
        #print(all_images)
        


    #for Speaker notes : 
   

    

    
    ppt=Presentation(pptfile_p)

    notes = []
    texty=[]

    for page, slide in enumerate(ppt.slides):
        # this is the notes that doesn't appear on the ppt slide,
        # but really the 'presenter' note. 
        textNote = slide.notes_slide.notes_text_frame.text
        notes.append((page,textNote))
        texty.append(textNote)
        

        #print(notes)
    #print(texty)
    cleanedtext=[]
    num = len(notes)
    for i in range(0,len(texty)):
        cleanedtext.append(re.sub('\x0b',' ',texty[i]))
        
    #print(cleanedtext)
    #for creating a doc
    
    document = Document()
    table = document.add_table(rows=num,cols=2)
    table.style = document.styles['TableGrid']
    table.style.font.name = 'Lora'
    table.style.font.size = Pt(20)
    
    
    for i in range(len(all_images)):
        cell= table.rows[i].cells[0]
        paragraph = cell.paragraphs[0]
        run = paragraph.add_run()
        #for loop for images
    
        run.add_picture(all_images[i],width=Inches(2.77),height =Inches(1.56))
        

        
    for j in range(0,len(cleanedtext)) :
        #print(texty[j])
        table.rows[j].cells[1].text = cleanedtext[j]
        #for loop for text
    
    
    document.save(pptfile_res+'\\NEW.docx')
    
    
    #filename_master is name of the file you want to merge the docx file into
    #master = Document_compose(temppath)
    master = Document_compose(pathy)

    composer = Composer(master)
    #filename_second_docx is the name of the second docx file
    doc2 = Document_compose(pptfile_res+'\\NEW.docx')
    #append the doc2 into the master using composer.append function
    composer.append(doc2)
    #Save the combined docx with a name
    composer.save(pptfile_res+"\\FINAL.docx")
    os.remove(pptfile_res+'\\NEW.docx')




def get_pdf():

    
    pptfile_p= str(my_str1.get())
    pptfile_res=str(my_str2.get())
    pptfile_p=pptfile_p.replace('/','\\')
    pptfile_res=pptfile_res.replace('/','\\')
    format_code = 17

    time_start = time.time()

    # create the MS word app
    word_app = comtypes.client.CreateObject('Word.Application')
    word_app.Visible = False

    file_input = os.path.abspath(pptfile_res+"\\FINAL.docx")
    file_output = os.path.abspath(pptfile_res+"\\FINAL.pdf")
    word_file = word_app.Documents.Open(file_input)
    word_file.SaveAs(file_output,FileFormat=format_code)
    word_file.Close()

    # close file and application
    word_app.Quit()

    time_end = time.time()

        
    
# the label for user_name

def browsefunc():
    filename = filedialog.askopenfilename()
    if (filename):
        my_str1.set(filename)
        
def browsefunc2():
    filename = filedialog.askdirectory()
    if (filename):
        my_str2.set(filename)

def browsefunc1():
    filename = filedialog.askopenfilename()
    if (filename):
        my_str3.set(filename)


my_str1=StringVar()
my_str2=StringVar()
my_str3=StringVar()
   
    
user_name = Label(top,
                  text = "Ppt file path", foreground = "green").place(x = 40,
                                           y = 60)
browsebutton = Button(top, text="Browse file location", command=browsefunc).place(x=450,y=60)

browsebutton = Button(top, text="Browse where to save file", command=browsefunc2).place(x=450,y=100)

#browsebutton = Button(top, text="Template", command=browsefunc1).place(x=380,y=140)

user_name2 = Label(top,
                  text = "Where To Save", foreground = "green").place(x = 40,
                                           y = 100) 

submit_button = Button(top,
                       text = "Get Word File",foreground = "green",command = give_path).place(x = 120,
                                              y = 200)

submit_button = Button(top,
                       text = "Get Pdf File",foreground = "green", command = get_pdf).place(x = 220, y = 200)



'''
user_name_input_area = Entry(top,
                             width = 30).place(x = 110,
                                               y = 60)
'''



path_label1 = Label(top,
                  textvariable=my_str1, foreground = "green").place(x = 150,
                                           y = 60) 
   
path_label2=  Label(top,
                  textvariable=my_str2, foreground = "green").place(x = 150,
                                           y = 100) 

   
path_label3=  Label(top,
                  textvariable=my_str3, foreground = "green").place(x = 150,
                                           y = 140) 
top.mainloop()
